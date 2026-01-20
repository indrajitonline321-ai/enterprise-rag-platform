from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from azure.storage.blob import BlobServiceClient
import pypdf
import os
from typing import List, Dict, Any
import io
import requests
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image
import pandas as pd
import os
import re
from typing import List
import requests
import ollama
from qdrant_client import QdrantClient,models
from qdrant_client.http.models import PointStruct, VectorParams, Distance, Filter, FieldCondition, MatchValue
from datetime import datetime, timezone
import uuid
from qdrant_client.models import PointStruct
import hashlib

app = FastAPI(title="RAG Service", version="0.1.0")
QDRANT_URL = "http://localhost:6333"
COLLECTION_NAME = "EnterPrise_chunks"
client = QdrantClient(QDRANT_URL)



class IngestRequest(BaseModel):
    document_id: str
    blob_url: str
    user_id: str

class Chunk(BaseModel):
    content: str
    document_id: str
    file_name: str
    page: int | None = None
    type: str = "text" 

class IngestResponse(BaseModel):
    document_id: str
    chunk_count: int
    chunks: List[Chunk]
    file_type: str
    userID: str

def chunk_text(text: str, max_memory_mb: int = 50) -> list:
    CHUNK_SIZE = 600  # Smaller = more pages
    chunks = []
    memory_limit = max_memory_mb * 1024 * 1024 / 0.78  # bytes
    
    for i in range(0, len(text), 500):  # Smaller step
        if len(chunks) * 0.78 > memory_limit:
            break  # Stop if memory limit hit
        chunk = text[i:i + CHUNK_SIZE].strip()
        if len(chunk) > 100:
            chunks.append(chunk)
    return chunks

def get_embedding(text: str) -> List[float]:
    """Ollama local embeddings (free, offline)"""
    response = ollama.embeddings(
        model="nomic-embed-text",  # 137-dim, fast
        prompt=text[:2048]  # Ollama limit
    )
    return response['embedding']

def init_collection():    
    if client.collection_exists(COLLECTION_NAME):
        print("✅ Collection exists - appending mode")
        return
    client.create_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=VectorParams(size=768, distance=Distance.COSINE)
    )
    print("✅ Collection created (first time)")

def detect_file_type(url: str) -> str:
    ext = url.lower().split('.')[-1]
    return {
        'pdf': 'PDF', 'docx': 'WORD', 'doc': 'WORD', 
        'xlsx': 'EXCEL', 'xls': 'EXCEL', 'pptx': 'PPT',
        'txt': 'TEXT', 'rtf': 'TEXT'
    }.get(ext, 'UNKNOWN')

def chunk_ocr_text(text: str, chunk_size: int = 800, min_len: int = 40) -> list:
    chunks = []
    for i in range(0, len(text), chunk_size):
        piece = text[i:i + chunk_size].strip()
        if len(piece) >= min_len:
            chunks.append(piece)
    return chunks

def normalize_text(t: str) -> List[str]:
    """Lowercase, remove emojis/punct, split into words."""
    t = t.lower()
    t = re.sub(r"[^\w\s]", " ", t)  # keep word chars + spaces
    return [w for w in t.split() if w]

def preprocess_question(question: str) -> str:
    if '|' not in question:
        return question  # no pipe → use as-is
    
    # Split into two parts
    left, right = question.split('|', 1)  # first | only
    left = left.strip()
    right = right.strip()
    # Ask Ollama to check similarity
    sim_prompt = f"""Compare these two questions for similarity (same topic/intent):
    
Q1: {left}
Q2: {right}

Return ONLY "MATCH" if they are about the same thing, else "NOMATCH". No explanation."""
    
    response = ollama.generate(
        model="llama3:latest",
        prompt=sim_prompt,
        options={"temperature": 0.1, "num_predict": 10}  # deterministic
    )
    ollama_reply = response['response'].strip().upper()
    if "NOMATCH" in ollama_reply:
        return right
    else:
        return f"{left} | {right}" 

def search_logic(query: str, limit: int = 5,doc_ids: list[str]=None) -> List[Dict]:
    query_embedding = get_embedding(query)
    search_filter = None
    if doc_ids:
        search_filter = models.Filter(
            must=[
                models.FieldCondition(
                    key="userID", 
                    match=models.MatchAny(any=doc_ids)
                )
            ]
        )
    results = client.query_points(
        collection_name=COLLECTION_NAME,
        query=query_embedding,
        limit=limit * 3,
        query_filter=search_filter,
        with_payload=True,
    )
    
    def rerank_score(item):
        content = item["content"].lower()
        query_words = query.lower().split()
        matches = sum(1 for word in query_words if word in content)
        return item["score"] * (1 + matches * 0.3)
    
    raw_results = []

    
    # ✅ SAFE: Iterate without unpacking
    for point_item in results.points:
        # Handle different formats
        if isinstance(point_item, tuple):
            if len(point_item) == 2:
                point_tuple, score = point_item
            else:
                point_tuple = point_item[0]
                score = point_item.score if hasattr(point_item[1], 'score') else 0.0
        else:
            # Single ScoredPoint object
            point_tuple = point_item
            score = getattr(point_item, 'score', 0.0)
        
        raw_results.append({
            "id": point_tuple.id,
            "score": float(score),
            "content": point_tuple.payload.get("content", ""),
            "file_name": point_tuple.payload.get("file_name", "unknown"),
            "document_id": point_tuple.payload.get("document_id"),
            "page": point_tuple.payload.get("page"),
            "chunk_index": point_tuple.payload.get("chunk_index")
        })
    
    top_results = sorted(raw_results, key=rerank_score, reverse=True)[:limit]
    return top_results

    

def is_similar_to_recent(chunks: list, new_text: str, window: int = 10, threshold: float = 0.7) -> bool:
  
    new_words = set(normalize_text(new_text)[:80])  # limit for speed
    if not new_words:
        return False

    for c in chunks[-window:]:
        existing_words = set(normalize_text(c.content)[:80])
        if not existing_words:
            continue
        inter = len(new_words & existing_words)
        union = len(new_words | existing_words)
        if union == 0:
            continue
        jaccard = inter / union
        if jaccard >= threshold:
            return True
    return False

@app.on_event("startup")
async def startup():
    init_collection()
    print("✅ Qdrant collection ready")


@app.get("/health")
async def health():
    return {"status": "ok"}



@app.post("/ingest", response_model=IngestResponse)
async def ingest(req: IngestRequest):
    try:
        # Get Azure connection string from env
        connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")

        existing = client.scroll(
        collection_name=COLLECTION_NAME,
        scroll_filter=Filter(
            must=[FieldCondition(key="file_name", match=MatchValue(value=req.blob_url.split('/')[-1]))]
        ),
        limit=1,
        )
    
        file_exists = len(existing[0]) > 0
    
        if file_exists:
        # 2. Delete old version (new upload = new version)
            print(f"🔄 Updating existing file: {req.blob_url.split('/')[-1]}")
            client.delete(
            collection_name=COLLECTION_NAME,
            points_selector={"must": [{"key": "file_name", "match": {"value": req.blob_url.split('/')[-1]}}]}
        )

        if not connection_string:
            raise HTTPException(status_code=500, detail="AZURE_STORAGE_CONNECTION_STRING not set")

        # Download PDF from Azure Blob
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)
        file_name = req.blob_url

        parts = req.blob_url.replace("https://", "").split("/")
        account_name = parts[0]
        container_name = parts[1]
        blob_name = parts[2]

      
    
        blob_client = blob_service_client.get_blob_client(
            container=container_name, 
            blob=blob_name
        )
        print(req.blob_url,container_name,file_name)
        # Extract container and blob name from URL
        if "blob.core.windows.net" not in req.blob_url:
            raise HTTPException(status_code=400, detail="Invalid blob URL")
                    #file_bytes = response.content
        file_type = detect_file_type(req.blob_url)
     

        all_chunks: List[Chunk] = []
        ocr_seen = set()
    
        print(f"Processing {file_type}: {req.blob_url}")

        pdf_bytes = blob_client.download_blob().readall()
    #PDF HANDLING (text + tables + OCR)
        if file_type == 'PDF':
            with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
                total_pages = len(pdf.pages)
                print(f"📄 TOTAL PAGES: {total_pages}")
        
                ocr_seen = set()  # for generic dedupe, not hardcoded

            for page_num, page in enumerate(pdf.pages, start=1):
                print(f"Processing Page {page_num}/{total_pages}...")

            # 1) TEXT – always first
                text = (page.extract_text() or "").strip()
                if len(text) > 500:
                    text_chunks = chunk_text(text)
                    for chunk_content in text_chunks:
                        all_chunks.append(Chunk(
                            content=chunk_content,
                            document_id=req.document_id,
                            file_name=file_name,
                            userID=req.user_id,
                            page=page_num,
                            type="text",
                        ))
                else:
                    try:
                        img = page.to_image(resolution=250).original
                        ocr_text = pytesseract.image_to_string(img, lang="eng+hin+dev").strip() 
                    # generic dedupe: no hardcoded strings
                        if len(ocr_text) >= 40 and len(ocr_text.split()) >= 5:
                            full_ocr_block = f"🖼️ IMAGE p{page_num}: {ocr_text}"
                            if is_similar_to_recent(all_chunks, full_ocr_block, window=10, threshold=0.7):
                                print(f"  🖼️ OCR p{page_num} skipped as duplicate")
                            else:
                                text_hash = hash(ocr_text[:150].lower())
                                if text_hash not in ocr_seen:
                                    ocr_seen.add(text_hash)
                                    print(f"  🖼️ OCR p{page_num}: {len(ocr_text)} chars")

                                    for chunk_content in chunk_ocr_text(full_ocr_block,
                                                    chunk_size=800,
                                                    min_len=40):
                                        all_chunks.append(Chunk(
                                            content=chunk_content,
                                        document_id=req.document_id,
                                        file_name=file_name,
                                        userID=req.user_id,
                                        page=page_num,
                                        type="image",
                                        ))
                    except Exception as e:
                            print(f"  ⚠️ OCR skip p{page_num}: {e}")

            # 2) TABLES – after text, same page
                tables = page.extract_tables() or []
                
                page_has_table = False
                clean_tables = []

                for table in tables:    
                # table is a list of rows; each row is a list of cells
                    if not table or len(table) < 2:
                        continue  # too small

    # remove empty rows
                    rows = []
                    for row in table:
                        cells = [str(c or "").strip() for c in row]
                        if any(cells):
                            rows.append(cells)

    # require at least 2 non-empty rows with multiple columns
                    if len(rows) >= 2 and max(len(r) for r in rows) >= 2:
                        page_has_table = True
                        clean_tables.append(rows)

                if page_has_table:
                    for table_num, table in enumerate(tables):
                        if table and len(table) > 1 and any(
                            cell for row in table[1:] for cell in row if cell
                        ):
                            table_lines = []
                            for row in table:
                                row_text = [str(cell or "").strip() for cell in row]
                                if any(row_text):
                                    table_lines.append(" | ".join(row_text))

                            if len(table_lines) > 1:
                                final_table = (
                                    f"📊 TABLE {table_num + 1} (p{page_num}):\n"
                                    + "\n".join(table_lines)
                                )
                                if is_similar_to_recent(all_chunks, final_table, window=8, threshold=0.8):
                                    continue
                                for chunk_content in chunk_text(final_table):
                                    all_chunks.append(Chunk(
                                    content=chunk_content,
                                    document_id=req.document_id,
                                    file_name=file_name,
                                    page=page_num,
                                    userID=req.user_id,
                                    type="table",
                                ))
                

    #WORD DOC (.docx)
        elif file_type == 'WORD':
            doc = docx.Document(io.BytesIO(pdf_bytes))
            text = "\n".join([para.text for para in doc.paragraphs])
            all_chunks.append(Chunk(
                                content=(chunk_text(text)),
                                document_id=req.document_id,
                                file_name=file_name,
                                page=page_num,
                                userID=req.user_id,
                                type="text"
                            ))
        
        # Word tables
            for table in doc.tables:
                table_text = ""
                for row in table.rows:
                    table_text += " | ".join([cell.text.strip() for cell in row.cells]) + "\n"   
                if is_similar_to_recent(all_chunks, table_text, window=8, threshold=0.8):
                    continue 
                all_chunks.append(Chunk(
                                content=(chunk_text(f"WORD TABLE:\n{table_text}")),
                                document_id=req.document_id,
                                file_name=file_name,
                                page=page_num,
                                userID=req.user_id,
                                type="table"
                            ))
    #EXCEL (.xlsx)
        elif file_type == 'EXCEL':
            df = pd.read_excel(io.BytesIO(pdf_bytes), sheet_name=None)
            for sheet_name, sheet_df in df.items():
                sheet_text = f"SHEET: {sheet_name}\n{sheet_df.to_string(index=False)}"
                all_chunks.append(Chunk(
                                content=(chunk_text(sheet_text)),
                                document_id=req.document_id,
                                file_name=file_name,
                                page=page_num,
                                userID=req.user_id,
                                type="text"
                            ))
    
    #POWERPOINT (.pptx)
        elif file_type == 'PPT':
            prs = Presentation(io.BytesIO(pdf_bytes))
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                all_chunks.append(Chunk(
                                content=(chunk_text(f"SLIDE {slide_num}:\n{slide_text}")),
                                document_id=req.document_id,
                                file_name=file_name,
                                page=page_num,
                                userID=req.user_id,
                                type="text"
                            ))
    
    #TEXT FILES
        elif file_type == 'TEXT':
            text = pdf_bytes.decode('utf-8', errors='ignore')
            all_chunks.append(Chunk(
                                content=(chunk_text(text)),
                                document_id=req.document_id,
                                file_name=file_name,
                                page=page_num,
                                userID=req.user_id,
                                type="text"
                            ))
    
        else:
            raise HTTPException(400, f"Unsupported file type: {file_type}")

        points = []
        for idx, chunk in enumerate(all_chunks):
            embedding = get_embedding(chunk.content)
            points.append(PointStruct(
                id= int.from_bytes(hashlib.sha256(f"{chunk.document_id}_{idx}".encode()).digest(), 'big') % (2**64),
                vector=embedding,
                payload={
                    "content": chunk.content,
                    "document_id": chunk.document_id,
                    "file_name": file_name,  # full filename
                    "page": chunk.page,
                    "type": chunk.type,
                    "chunk_index": idx,
                    "userID" : req.user_id,
                    "uploaded_at": datetime.now(timezone.utc).isoformat(),  # add timestamp
                }
            ))
    
        client.upsert(collection_name=COLLECTION_NAME, points=points)
    
        action = "updated" if file_exists else "created"
        print(f"✅ {action}: {len(all_chunks)} vectors for {file_name}")

        return IngestResponse(
            document_id=req.document_id,
            chunk_count=len(all_chunks),
            file_type=file_type,
            chunks=all_chunks,
            userID=req.user_id
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    



class SearchRequest(BaseModel):
    query: str
    limit: int = 5

class ChatRequest(BaseModel):
    query: str
    limit: int = 5
    docIds: List[str]
    userId: str
    


@app.post("/search")
async def search(req: SearchRequest):
 
    return {"results": search_logic(req.query, req.limit)}

# Chat calls SAME logic automatically
@app.post("/chat")
async def chat(req: ChatRequest):
    question = preprocess_question(req.query) 
    # ✅ Automatically searches top 3
    top_chunks = search_logic(question, limit=3,doc_ids=req.docIds)
    context = "\n\n".join([r["content"] for r in top_chunks])
    prompt = f"""Using ONLY this context, answer:
CONTEXT: {context}

Q: {question}
A:"""
    
    response = ollama.chat(model="llama3:latest", messages=[{"role": "user", "content": prompt}])
    return {
        "answer": response['message']['content'],
        "sources": top_chunks  # Bonus: show sources!
    }
